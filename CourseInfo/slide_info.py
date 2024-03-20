import pandas as pd
from pptx import Presentation

class Pres():
    def __init__(self, module):
        self.module = int(module)
        # self.extract_info()

    def open_slide_deck(self, file):
        self.prs = Presentation(file)
        return(len(self.prs.slides))

    def extract_info(self):        
        data = pd.DataFrame(columns = ["Module", "Slide", "Title", "Notes", "Knowledge Application", "Discussion Comments and Questions", "Reference"])
        for page, slide in enumerate(self.prs.slides):
            page +=1
            try:
                title = slide.shapes.title.text
            except:
                title = pd.NA
            textNote = slide.notes_slide.notes_text_frame.text
            #STRUCTURE NOTES
            notes,knowApp,discuss,ref = self.structure_notes(textNote)
            data.loc[len(data)] = [self.module, page, title, notes, knowApp, discuss, ref]
        self.info = data

    def structure_notes(self, textNote):
        notes = textNote.split("Knowledge Application")[0]
        try:
            less_notes = textNote.split("Knowledge Application")[1]
            knowApp = less_notes.split("Discussion Comments and Questions")[0]
        except:
            knowApp = "Not applicable."
        try:
            less_knowapp = less_notes.split("Discussion Comments and Questions")[1]
            discuss = less_knowapp.split("References")[0]
        except:
            discuss = "Not applicable."
        try:
            ref = less_knowapp.split("References")[1]
        except:
            ref = "Not applicable."
        return(notes, knowApp, discuss,ref)
        
    # def add_notes():
    #     pass